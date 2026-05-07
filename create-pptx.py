from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 format
prs.slide_height = Inches(7.5)

# Bo'sh slayd maketlari
title_slide_layout = prs.slide_layouts[0]   # Title Slide (faqat sarlavha)
content_slide_layout = prs.slide_layouts[1] # Title and Content (sarlavha + matn)
title_only_layout = prs.slide_layouts[5]    # Title Only (sarlavha va bo'sh maydon)

def add_bullet_slide(title, bullet_list):
    """Sarlavha va markerli ro'yxat yaratish"""
    slide = prs.slides.add_slide(content_slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for item in bullet_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
    return slide

def add_table_slide(title, headers, rows):
    """Jadval ko'rinishidagi slayd yaratish"""
    slide = prs.slides.add_slide(title_only_layout)
    slide.shapes.title.text = title
    # Jadval o'lchamlari
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(11.3)
    height = Inches(0.8 * (len(rows) + 1))  # taxminiy balandlik
    table_shape = slide.shapes.add_table(len(rows)+1, len(headers), left, top, width, height)
    table = table_shape.table

    # Sarlavha satri
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(18)

    # Ma'lumot satrlari
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx+1, col_idx)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
    return slide

# ===================== SLIDE 1: SARLAVHA =====================
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
title.text = ("Aholiga shoshilinch, kechiktirib bo'lmaydigan va tez tibbiy yordamni tashkil qilishning\n"
              "o'ziga xos xususiyatlari. Tez tibbiy yordam stansiyalarining ish faoliyati va ularni baholash ko'rsatkichlari")
title.text_frame.paragraphs[0].font.size = Pt(32)

# ===================== SLIDE 2: Asosiy tushunchalar =====================
add_bullet_slide("Asosiy tushunchalar", [
    "Shoshilinch tibbiy yordam – bemor hayotiga bevosita xavf soluvchi to'satdan paydo bo'lgan holatlarda ko'rsatiladigan yordam (masalan, infarkt, og'ir jarohatlar).",
    "Kechiktirib bo'lmaydigan tibbiy yordam – hayot uchun bevosita xavf bo'lmasa-da, qisqa vaqt ichida ko'rsatilmasa, sog'liq uchun jiddiy asoratlarga olib kelishi mumkin bo'lgan yordam (masalan, yuqori haroratli tutilishlar, o'tkir og'riq sindromi).",
    "Tez tibbiy yordam (TTY) – yuqoridagi ikki holatni o'z ichiga olgan, aholiga kecha-yu kunduz, bepul va chaqiriq asosida ko'rsatiladigan maxsus tibbiy xizmat turi."
])

# ===================== SLIDE 3: Tashkil qilish xususiyatlari 1/3 =====================
add_bullet_slide("Tez tibbiy yordamni tashkil qilishning o'ziga xos xususiyatlari (1/3)", [
    "Kechayu kunduz ishlash rejimi: TTY xizmati 24/7, bayram va dam olish kunlarisiz uzluksiz faoliyat yuritadi.",
    "Tezkorlik: Chaqiriq qabul qilingandan so'ng ekipajning yetib borish vaqti qat'iy normativlar bilan chegaralangan (odatda shaharda 15-20 daqiqagacha).",
    "Mobil va statsionar bosqich: Yordam ko'rsatish joyida (bemor yonida), transportirovka vaqtida va stasionarga yetkazilgach davom etadi."
])

# ===================== SLIDE 4: Tashkil qilish xususiyatlari 2/3 =====================
add_bullet_slide("Tez tibbiy yordamni tashkil qilishning o'ziga xos xususiyatlari (2/3)", [
    "Birlamchi tashxis va saralash: TTY shifokori (yoki feldsher) dastlabki tashxis qo'yadi, bemorning og'irlik darajasini aniqlab, uni mos davolash muassasasiga yo'naltiradi.",
    "Multidisiplinar yondashuv: Zarur hollarda bir vaqtning o'zida bir nechta profildagi brigadalar (kardiologik, reanimatsion, pediatrik) jalb qilinadi.",
    "Bemor oqimini boshqarish: Shoshilinch yordam xonalariga tushum bosimini kamaytirish uchun chaqiriqlarni asosli saralash va konsultativ yordam markazlari bilan bog'lash."
])

# ===================== SLIDE 5: Tashkil qilish xususiyatlari 3/3 =====================
add_bullet_slide("Tez tibbiy yordamni tashkil qilishning o'ziga xos xususiyatlari (3/3)", [
    "Normativ-huquqiy baza: Faoliyat Sog'liqni saqlash vazirligi buyruqlari, standartlar va klinik protokollar asosida tartibga solinadi.",
    "Moddiy-texnik ta'minot: Avtotransport, zamonaviy tibbiy asbob-uskunalar, dorilar va aloqa vositalari bilan jihozlanganlik darajasi yuqori bo'lishi shart.",
    "Ma'lumotlar bazasi: Chaqiriqlar yagona dispetcherlik markazi orqali qabul qilinadi, har bir holat elektron kartada qayd etilib, statistik tahlil yuritiladi."
])

# ===================== SLIDE 6: Stansiya tuzilmasi =====================
add_bullet_slide("Tez tibbiy yordam stansiyalarining tuzilmasi", [
    "Operativ dispetcherlik bo'limi: 103 chaqiriqlarini qabul qiladi, saralaydi va brigadalarni muvofiqlashtiradi.",
    "Statsionar bo'lim (kasalxona): Ba'zi yirik stansiyalarda qisqa muddatli kuzatuv palatalari mavjud.",
    "Ko'chma brigadalar: Vrach brigadalari (umumiy va ixtisoslashgan), feldsher brigadalari, reanimatsion brigadalar.",
    "Maslahat markazi (konsultativ yordam): Masofaviy EKG uzatish, toksikologik maslahat va boshqa yo'nalishlar."
])

# ===================== SLIDE 7: Asosiy vazifalar =====================
add_bullet_slide("TTY stansiyasining asosiy vazifalari", [
    "Aholiga kechiktirib bo'lmaydigan va shoshilinch tibbiy yordam ko'rsatish.",
    "Bemorlarni tibbiy muassasalarga qat'iy ko'rsatmalar asosida tashish (reanimatsiya, tug'ruq, jarrohlik markazlari).",
    "Favqulodda vaziyatlarda, ommaviy halokatlarda yordam tashkil qilish.",
    "Kasalxonalardan tashqarida o'lim holatlarini qayd etish va tegishli xulosalar berish.",
    "Aholi o'rtasida birinchi yordam ko'rsatishni targ'ib qilishda ishtirok etish."
])

# ===================== SLIDE 8: Baholash – Tezkorlik (JADVAL) =====================
add_table_slide("TTY faoliyatini baholash ko'rsatkichlari: Tezkorlik",
                ["Ko'rsatkich", "Tavsifi"],
                [
                    ["Chaqiriqqa yetib borish vaqti",
                     "Chaqiriq qabul qilingandan brigada manzilga yetguncha o'tgan vaqt. Normativ: 15-20 daqiqa (shahar)."],
                    ["Chaqiriqqa xizmat ko'rsatish umumiy vaqti",
                     "Chaqiriq qabulidan brigadaning navbatchi stansiyaga qaytgunigacha yoki yangi topshiriq olguncha o'tgan vaqt."],
                    ["Kechikishlar ulushi",
                     "Normativdan ortiq vaqtda yetib borilgan chaqiriqlar foizi."],
                    ["Dispetcher tomonidan qabul qilish tezligi",
                     "Qo'ng'iroqni javobsiz kutish vaqti (sekundlarda)."]
                ])

# ===================== SLIDE 9: Baholash – Sifat va natijadorlik =====================
add_bullet_slide("TTY faoliyatini baholash ko'rsatkichlari: Sifat va natijadorlik", [
    "Takroriy chaqiriqlar foizi: Birinchi chaqiriqdan so'ng 24 soat ichida o'sha bemorga qayta chaqiruv kelishi. Yuqori ko'rsatkich noto'g'ri tashxis yoki yetarli yordam ko'rsatilmaganini bildiradi.",
    "Tashxislarning moslik darajasi: TTY shifokori qo'ygan dastlabki tashxis bilan kasalxonadagi yakuniy tashxis o'rtasidagi muvofiqlik.",
    "O'lim holatlari tahlili: Brigada yetib kelguncha va brigada ishtirokida sodir bo'lgan o'limlar, ularning oldini olish mumkinligi bo'yicha ekspert bahosi.",
    "Asoratlarsiz muvaffaqiyatli yordam: Joyida samarali yordam ko'rsatilib, bemor og'ir ahvoldan chiqarilgan holatlar soni (muvaffaqiyatli reanimatsiya ulushi)."
])

# ===================== SLIDE 10: Baholash – Resurs va ish yuki =====================
add_bullet_slide("TTY faoliyatini baholash ko'rsatkichlari: Resurs va ish yuki", [
    "Aholining 1000 nafariga to'g'ri keladigan chaqiriqlar soni: TTY xizmatiga bo'lgan talab va yuklamani ko'rsatadi.",
    "Brigadaning o'rtacha kunlik yuklamasi: Bir brigada tomonidan sutkada bajarilgan chaqiriqlar o'rtacha soni.",
    "Kadrlar bilan ta'minlanganlik: Shifokor va feldsher lavozimlarining to'ldirilish foizi, kadrlar qo'nimsizligi.",
    "Avtotransportning eskirish ko'rsatkichi va texnik tayyorgarlik koeffitsiyenti: Mashinalarning doimiy safarbarlik holatini ifodalaydi."
])

# ===================== SLIDE 11: Muammolar va yechimlar (JADVAL) =====================
add_table_slide("Asosiy muammolar va takomillashtirish yo'llari",
                ["Muammo", "Yechim"],
                [
                    ["Asossiz chaqiriqlarning ko'pligi",
                     "Aholi o'rtasida sanitariya marifati ishlari, triaj tizimini kuchaytirish"],
                    ["Brigadalarning kechikib borishi",
                     "GPS-navigatsiya, yo'l harakati bilan bog'liq imtiyozlar, qo'shimcha postlar ochish"],
                    ["Moliyaviy va moddiy ta'minot",
                     "Maqsadli davlat dasturlari, pullik xizmat turlarini kengaytirish"],
                    ["Kadrlar yetishmovchiligi",
                     "Moddiy va ma'naviy rag'batlantirish, ta'lim grantlarini ko'paytirish"]
                ])

# ===================== SLIDE 12: Xulosa =====================
add_bullet_slide("Xulosa", [
    "Tez tibbiy yordam – sog'liqni saqlash tizimining eng muhim va tezkor bo'g'ini.",
    "Uning o'ziga xosligi – kechayu kunduzlik, kechiktirib bo'lmaslik va yuqori malaka talab qilishdadir.",
    "Stansiyalar faoliyatini baholash uchun faqat miqdoriy ko'rsatkichlar emas, balki davolash sifati va tezkorligi asosiy mezon hisoblanadi.",
    "Zamonaviy texnologiyalarni joriy etish va doimiy tahliliy monitoring xizmat samaradorligini oshirishning bosh omilidir."
])

# Slaydlarni saqlash
output_path = "Tez_tibbiy_yordam_slaydlar.pptx"
prs.save(output_path)
print(f"Taqdimot muvaffaqiyatli yaratildi: {os.path.abspath(output_path)}")